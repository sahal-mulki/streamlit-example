from openai import OpenAI
import io
import json
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches


def create_layout(n_slides: int, audience: str, topic: str, openaikey: str):
  # Create layout for PPT

  client = OpenAI(api_key=openaikey)

  response = client.chat.completions.create(
    model="gpt-3.5-turbo-1106",
    response_format={ "type": "json_object" },
    messages=[
      {
        "role": "system",
        "content": "Make a layout for a powerpoint. Be sure to make it engaging and use emojis. Use some slides for questions and tasks. It should be in this format:\n{slides: {\n            {\"id\": \"slide1\", \"title\": \"Introduction\", \"content\": \"....\"},\n              {\"id\": \"slide2\", \"title\": \"Starter\", \"content\": \"......\"},\n              {\"id\": \"slide3\", \"title\": \"Explanation\", \"content\": ........\"}\n}\n            ADHERE TO THE JSON FORMAT AND DO NOT PROVIDE ANY ADDITIONAL INFORMATION. Do not use any controversial or political statements in the presentation.         "
      },
      {
        "role": "user",
        "content": "Topic: " + topic + " , Audience: " + audience + ", Number of Slides: " + str(n_slides)
      }
    ],
    temperature=0.3,
    max_tokens=597,
    top_p=1,
    frequency_penalty=0.35,
    presence_penalty=0.11
  )

  exec("slides = " + response.choices[0].message.content, globals())

  slides_list_title = []

  try:
    for x in range(1, n_slides+1):
      slides_list_title.append({"id": str("slide-" + str(x)), "title": slides["slides"][str("slide" + str(x))]["title"], "content": slides["slides"][str("slide" + str(x))]["content"]})
  except:
    client = OpenAI(api_key=openaikey)

    response = client.chat.completions.create(
      model="gpt-3.5-turbo-1106",
      response_format={ "type": "json_object" },
      messages=[
        {
          "role": "system",
          "content": "Make a layout for a powerpoint. Be sure to make it engaging and use emojis. Use some slides for questions and tasks. It should be in this format:\n{slides: {\n            {\"id\": \"slide1\", \"title\": \"Introduction\", \"content\": \"....\"},\n              {\"id\": \"slide2\", \"title\": \"Starter\", \"content\": \"......\"},\n              {\"id\": \"slide3\", \"title\": \"Explanation\", \"content\": ........\"}\n}\n            ADHERE TO THE JSON FORMAT AND DO NOT PROVIDE ANY ADDITIONAL INFORMATION. Do not use any controversial or political statements in the presentation.         "
        },
        {
          "role": "user",
          "content": "Topic: " + topic + " , Audience: " + audience + ", Number of Slides: " + str(n_slides)
        }
      ],
      temperature=0.3,
      max_tokens=597,
      top_p=1,
      frequency_penalty=0.35,
      presence_penalty=0.11
    )

  
    for x in range(1, n_slides+1):
      slides_list_title.append({"id": str("slide-" + str(x)), "title": slides["slides"][str("slide" + str(x))]["title"], "content": slides["slides"][str("slide" + str(x))]["content"]})

  return slides_list_title

def create_content(slides_list_titles: dict, audience: str, openaikey: str):

  client = OpenAI(api_key=openaikey)
  
  response = client.chat.completions.create(
      model="gpt-3.5-turbo-1106",
      response_format={ "type": "json_object" },
      messages=[
        {
          "role": "system",
          "content": """"Give me text to put in a presentation slide.  Give me one large-sized paragraph per slide. Be sure to make it engaging and use emojis. Use some slides for questions and tasks. Provide your answer in this format:
                  {content: "................................"}
                  ADHERE TO THE JSON FORMAT AND DO NOT PROVIDE ANY ADDITIONAL INFORMATION. Do not use any controversial or political statements in the presentation. Give me one large-sized paragraph per slide.
              """
        },
        {
          "role": "user",
          "content": "Slide Title: " + slides_list_titles["title"] + " , Audience: " + audience + ", Slide Description: " + slides_list_titles["content"]
        }
      ],
      temperature=0.7,
      max_tokens=1200,
      top_p=1,
      frequency_penalty=0.35,
      presence_penalty=0.11
  )

  return response.choices[0].message.content

def create_font_list(topic: str, audience: str, openaikey: str):

  client = OpenAI(api_key=openaikey)

  response = client.chat.completions.create(
      model="gpt-3.5-turbo-1106",
      response_format={ "type": "json_object" },
      messages=[
        {
          "role": "system",
          "content": """What font would be appropriate for a presentation? Provide font that are included in Microsoft Powerpoint.
              PROVIDE NO OTHER INFORMATION THAN THE BELOW IN JSON FORMAT. USE ENGAGING AND BEAUTIFUL FONTS.
              Provide your answer as such:
              {"fonts": [font, font, font, font]}
              Do not provide any other information than the font in THE GIVEN format. DO NOT ADD ANY COMMENTS OR EXTRA INFORMATION OR AN EXPLANATION OF YOUR CHOICES. If you are unsure of any font's existance, do not include it in the list. DO NOT GIVE ANY MORE INFORMATION THAN THE LIST.
              """
        },
        {
          "role": "user",
          "content": "Topic: " + topic + " , Audience: " + audience
        }
      ],
      temperature=0.3,
      max_tokens=597,
      top_p=1,
      frequency_penalty=0.35,
      presence_penalty=0.11
  )

  exec("fonts = " + response.choices[0].message.content, globals())

  return fonts



def format_slides(slides: dict, slides_detailed: list):
  slides_final = []

  pee = -1

  for x in slides:

    pee += 1

    slides_final.append({"title": slides[pee]["title"], "content": slides_detailed[pee]["content"]})
  
  return slides_final



def create_ppt(slides_content: list, template_choice: int, presentation_title: str, fonts: list):

    def listToString(s):

      # initialize an empty string
      str1 = ""

      # traverse in the string
      for ele in s:
          str1 += str(ele + "\n")

      # return string
      return str1


    if template_choice == 1:
      template_choice = "Design-2"
    elif template_choice == 2:
      template_choice = "Design-4"
    elif template_choice == 3:
      template_choice = "Design-6"
    elif template_choice == 4:
      template_choice = "Design-7"
    else:
      template_choice = "Design-2"


    template_path = os.path.join("ppt-templates/", f"{template_choice}.pptx")


    prs = Presentation(template_path)

    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    # add title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = presentation_title

    if template_choice == 'dark_modern':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = fonts[0]
                run.font.color.rgb = RGBColor(255, 165, 0)  # RGB for orange color

    else:
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = fonts[0]
                run.font.color.rgb = RGBColor(255, 20, 147)  # RGB for deep pink color

    pooooo = -1

    # add content slides
    for slide_content in slides_content:
        pooooo = pooooo + 1
        slide = prs.slides.add_slide(content_slide_layout)

        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = slide_content['title']
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if template_choice == "Organisch" or "Badge" or "Zitierfahig":
                          run.font.name = fonts[0]
                          run.font.color.rgb = RGBColor(255, 165, 0)  # RGB for orange color
            elif placeholder.placeholder_format.type == 7:  # Content
                try:
                  llll = slides_multimedias[pooooo]['image_description']
                  placeholder.text = slide_content['content']
                  for paragraph in placeholder.text_frame.paragraphs:
                      #shapes = slide.shapes
                      #body_shape = shapes.placeholders[1]
                      placeholder.width = Inches(7)
                      placeholder.height = Inches(12)
                      for run in paragraph.runs:
                          if template_choice == "Design-4":
                            run.font.name = fonts[0]
                            run.font.color.rgb = RGBColor(255, 255, 255)  # RGB for black color
                          else:
                            run.font.name = fonts[0]
                            run.font.color.rgb = RGBColor(0, 0, 0)  # RGB for white color
                except:
                  #x = slide_content['content'].split("\n")
                  #x = x[:7]
                  placeholder.text = slide_content['content']
                  for paragraph in placeholder.text_frame.paragraphs:
                      for run in paragraph.runs:
                          if template_choice == "Design-4":
                            run.font.name = fonts[0]
                            run.font.color.rgb = RGBColor(255, 255, 255)  # RGB for black color
                          else:
                            run.font.name = fonts[0]
                            run.font.color.rgb = RGBColor(0, 0, 0)  # RGB for white color

    ## add credits slide

    # Delete the first two slides after all new slides have been added
    #delete_first_two_slides(prs)

    # Save the presentation
    prs.save(os.path.join('generated-ppt/', 'generated_presentation.pptx'))
  